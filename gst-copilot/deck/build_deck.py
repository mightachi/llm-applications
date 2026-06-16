"""Generate GST_Copilot_Pitch.pptx from pitch_deck.md.

Parses the Markdown (one slide per '## Slide' heading), rendering bullets, simple
pipe-tables, and speaker notes (the 'Notes:' block). Run:  python deck/build_deck.py
"""

from __future__ import annotations

import os
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

HERE = os.path.dirname(__file__)
SRC = os.path.join(HERE, "pitch_deck.md")
OUT = os.path.join(HERE, "GST_Copilot_Pitch.pptx")

INK = RGBColor(0x1E, 0x29, 0x3B)
ACCENT = RGBColor(0x4F, 0x46, 0xE5)
MUTED = RGBColor(0x64, 0x74, 0x8B)


def parse_slides(md: str) -> list[dict]:
    slides = []
    for block in re.split(r"^##\s+Slide.*$", md, flags=re.M)[1:]:
        block = block.strip()
        notes = ""
        if "Notes:" in block:
            block, notes = block.split("Notes:", 1)
        lines = [ln.rstrip() for ln in block.strip().splitlines()]
        # First non-empty content line is the title (strip markdown bold/marks).
        title = ""
        body_lines = []
        for ln in lines:
            if not title and ln.strip() and not ln.startswith("|"):
                title = re.sub(r"[*_#]", "", ln).strip()
                continue
            body_lines.append(ln)
        slides.append({"title": title, "body": body_lines, "notes": notes.strip()})
    return slides


def _add_bullets(tf, lines):
    first = True
    for ln in lines:
        text = ln.strip()
        if not text:
            continue
        text = re.sub(r"^[-*]\s+", "", text)
        text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = text
        p.font.size = Pt(16)
        p.font.color.rgb = INK
        p.level = 0


def _add_table(slide, rows):
    n_rows, n_cols = len(rows), len(rows[0])
    gfx = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.6), Inches(2.0), Inches(9.0), Inches(0.4 * n_rows)
    )
    table = gfx.table
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = re.sub(r"\*\*(.+?)\*\*", r"\1", val.strip())
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = INK if r else RGBColor(0xFF, 0xFF, 0xFF)


def _parse_pipe_table(lines):
    rows = []
    for ln in lines:
        if ln.strip().startswith("|"):
            cells = [c.strip() for c in ln.strip().strip("|").split("|")]
            if all(set(c) <= {"-", ":", " "} for c in cells):
                continue  # separator row
            rows.append(cells)
    return rows


def build():
    with open(SRC, encoding="utf-8") as f:
        md = f.read()
    slides = parse_slides(md)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for s in slides:
        slide = prs.slides.add_slide(blank)
        # Title bar
        box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9.0), Inches(1.0))
        tf = box.text_frame
        tf.word_wrap = True
        tf.text = s["title"]
        tf.paragraphs[0].font.size = Pt(30)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = ACCENT

        table_rows = _parse_pipe_table(s["body"])
        if table_rows:
            _add_table(slide, table_rows)
        else:
            body = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(9.0), Inches(5.0))
            body.text_frame.word_wrap = True
            _add_bullets(body.text_frame, s["body"])

        if s["notes"]:
            slide.notes_slide.notes_text_frame.text = s["notes"]

    prs.save(OUT)
    print(f"wrote {OUT} ({len(slides)} slides)")


if __name__ == "__main__":
    build()
